#!/usr/bin/env python3
"""Build the founder-framed capital-map deck (python-pptx, no Node needed)."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---- palette (from the dashboard's own CSS vars) ----
BG    = "0E0E13"
CARD  = "181820"
INK   = "F1F1F5"
MUTE  = "A6A6B2"
FAINT = "73737E"
AI    = "FF7A18"
WEB3  = "77E26D"
FIN   = "B88CFF"
TEAL  = "2DD4BF"
LINE  = "2A2A33"

HEAD = "Cambria"   # safe serif for headers
BODY = "Calibri"   # safe sans for body

def C(h): return RGBColor.from_string(h)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

def slide():
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = C(BG)
    return s

def card(s, x, y, w, h, fill=CARD, radius=0.06, line=None):
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = C(fill)
    if line:
        shp.line.color.rgb = C(line); shp.line.width = Pt(1)
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    try: shp.adjustments[0] = radius
    except Exception: pass
    return shp

def rect(s, x, y, w, h, fill):
    shp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = C(fill)
    shp.line.fill.background(); shp.shadow.inherit = False
    return shp

def text(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space=1.0, para_after=0.0):
    """runs: list of paragraphs; each paragraph is list of (txt, size, color, bold, italic, font)."""
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = space
        if para_after: p.space_after = Pt(para_after)
        for (txt, size, color, bold, italic, font) in para:
            r = p.add_run(); r.text = txt
            r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
            r.font.name = font; r.font.color.rgb = C(color)
    return tb

def R(txt, size, color=INK, bold=False, italic=False, font=BODY):
    return (txt, size, color, bold, italic, font)

def kicker(s, txt, color=FAINT):
    text(s, 0.7, 0.45, 11.0, 0.4, [[R(txt, 12, color, True, False, BODY)]])

def vbar(s, x, base_y, w, val, vmax, area_h, color, val_lbl, yr_lbl):
    h = max(0.05, (val / vmax) * area_h)
    rect(s, x, base_y - h, w, h, color)
    text(s, x - 0.25, base_y - h - 0.42, w + 0.5, 0.4, [[R(val_lbl, 13, INK, True)]], align=PP_ALIGN.CENTER)
    text(s, x - 0.25, base_y + 0.06, w + 0.5, 0.32, [[R(yr_lbl, 12, MUTE)]], align=PP_ALIGN.CENTER)

# ============================================================ SLIDE 1 — title
s = slide()
kicker(s, "VENTURE CAPITAL  ·  2022–2025  ·  EVERY FIGURE SOURCED", AI)
text(s, 0.7, 2.05, 11.6, 2.4, [
    [R("A founder's read on the ", 40, INK, True, False, HEAD),
     R("capital map.", 40, AI, True, False, HEAD)],
    [R("Where venture money actually went — and the questions it raises", 21, MUTE)],
    [R("if you're building in crypto or fintech.", 21, MUTE)],
], space=1.06, para_after=6)
text(s, 0.7, 6.5, 12.0, 0.6, [[
    R("Backbone: Crunchbase annual · Web3 deep-dive: Galaxy Digital quarterly · ", 12, FAINT),
    R("gitlostinthesauce.github.io/InsubordinateandChurlish_VCflow", 12, TEAL),
]])

# ============================================================ SLIDE 2 — the gravity
s = slide()
kicker(s, "01  ·  THE GRAVITY", AI)
text(s, 0.7, 0.78, 12.0, 0.7, [[R("AI is where the capital — and the attention — is", 30, INK, True, False, HEAD)]])
text(s, 0.7, 1.55, 12.0, 0.5, [[R("The question: ", 16, FAINT, True),
    R("if I'm not an AI company, am I invisible to capital?", 16, INK, False, True)]])

# left: AI ascent bars (real annual values, $B)
text(s, 0.7, 2.35, 5.4, 0.35, [[R("AI venture funding, $B per year", 13, MUTE, True)]])
vals = [("2022", 45.8, "$45.8B"), ("2023", 50.0, "$50B"), ("2024", 114.0, "$114B"), ("2025", 211.0, "$211B")]
base_y, area_h, bw, x0, gap = 6.35, 2.75, 0.95, 1.05, 1.35
for i, (yr, v, lbl) in enumerate(vals):
    vbar(s, x0 + i * gap, base_y, bw, v, 211.0, area_h, AI, lbl, yr)

# right: three stat cards
stats = [
    ("~50%", "of ALL global venture dollars went to AI in 2025", AI),
    ("4.6×", "AI's run-up in three years — $45.8B (2022) to $211B (2025)", INK),
    ("~$93B", "of that landed in just three labs: OpenAI, Anthropic, xAI", AI),
]
cx, cw = 6.95, 5.7
cy, ch, cgap = 2.3, 1.28, 0.18
for i, (big, lbl, col) in enumerate(stats):
    y = cy + i * (ch + cgap)
    card(s, cx, y, cw, ch)
    text(s, cx + 0.35, y, 2.1, ch, [[R(big, 34, col, True, False, HEAD)]], anchor=MSO_ANCHOR.MIDDLE)
    text(s, cx + 2.55, y, cw - 2.85, ch, [[R(lbl, 14, INK)]], anchor=MSO_ANCHOR.MIDDLE, space=1.02)

text(s, 0.7, 6.85, 12.0, 0.5, [[R("Read:  ", 14, WEB3, True),
    R("the narrative gravity is overwhelming — but the AI core is brutally concentrated and expensive to enter.", 14, MUTE, False, True)]])

# ============================================================ SLIDE 3 — your own sector
s = slide()
kicker(s, "02  ·  YOUR OWN LANE", AI)
text(s, 0.7, 0.78, 12.0, 0.7, [[R("Crypto and fintech aren't dead — they're disciplined", 30, INK, True, False, HEAD)]])
text(s, 0.7, 1.55, 12.0, 0.5, [[R("The question: ", 16, FAINT, True),
    R("is my market actually fundable right now?", 16, INK, False, True)]])

def sector_card(x, color, title, rows):
    cy2, cw2, ch2 = 2.25, 5.75, 3.95
    card(s, x, cy2, cw2, ch2)
    text(s, x + 0.45, cy2 + 0.35, cw2 - 0.9, 0.5, [[R(title, 20, color, True, False, HEAD)]])
    runs = []
    for r in rows:
        runs.append([R("›  ", 15, color, True), R(r, 15, INK)])
    text(s, x + 0.45, cy2 + 1.05, cw2 - 0.9, ch2 - 1.4, runs, space=1.04, para_after=9)

sector_card(0.7, WEB3, "Crypto / Web3", [
    "Collapsed −86% into the 2023 winter ($14.2B → $1.93B quarterly).",
    "Back to $8.5B in Q4 2025 — strongest quarter since Q2 2022.",
    "The money returned to infrastructure & trading rails, not consumer/NFT froth.",
    "One $2B Binance/MGX deal was ~41% of a single quarter.",
])
sector_card(6.85, FIN, "Fintech", [
    "Halved from $90.2B (2022) to $43B (2023) — a 6-year low.",
    "Rebuilt to $51.8B (2025), +27% YoY.",
    "Recovery shaped by fewer, larger rounds — not a return to spray.",
    "Stripe's $6.5B round was ~15% of all fintech funding in 2023.",
])

text(s, 0.7, 6.5, 12.0, 0.5, [[R("Read:  ", 14, WEB3, True),
    R("capital is back in both lanes — but it rewards focus and category leadership, not breadth.", 14, MUTE, False, True)]])

# ============================================================ SLIDE 4 — the pattern
s = slide()
kicker(s, "03  ·  THE PATTERN ACROSS BOTH", AI)
text(s, 0.7, 0.78, 12.0, 0.7, [[R("Winners take most — in every sector, not just AI", 30, INK, True, False, HEAD)]])
text(s, 0.7, 1.55, 12.0, 0.5, [[R("The question: ", 16, FAINT, True),
    R("what actually makes a company fundable in this market?", 16, INK, False, True)]])

# left: megaround growth, two bars + callout to their right (clear of bars)
text(s, 0.7, 2.35, 5.6, 0.35, [[R("Average landmark round size", 13, MUTE, True)]])
base_y, area_h = 6.35, 2.55
vbar(s, 1.45, base_y, 1.15, 0.44, 2.6, area_h, FAINT, "~$440M", "2022")
vbar(s, 3.15, base_y, 1.15, 2.6, 2.6, area_h, AI, "~$2.6B", "2025")
text(s, 4.8, 3.5, 1.9, 0.85, [[R("~6×", 44, AI, True, False, HEAD)]], align=PP_ALIGN.LEFT)
text(s, 4.83, 4.45, 2.0, 0.6, [[R("bigger in", 13, MUTE)], [R("three years", 13, MUTE)]], space=1.0)

# right: three concentration facts
facts = [
    ("AI", AI, "~half of all landmark capital → 3 labs"),
    ("Web3", WEB3, "~41% of a quarter → one $2B deal"),
    ("Fintech", FIN, "~15% of a year → Stripe's single round"),
]
cx, cw = 6.95, 5.7
cy, ch, cgap = 2.3, 1.28, 0.18
for i, (tag, col, lbl) in enumerate(facts):
    y = cy + i * (ch + cgap)
    card(s, cx, y, cw, ch)
    text(s, cx + 0.35, y, 1.9, ch, [[R(tag, 22, col, True, False, HEAD)]], anchor=MSO_ANCHOR.MIDDLE)
    text(s, cx + 2.3, y, cw - 2.6, ch, [[R(lbl, 14.5, INK)]], anchor=MSO_ANCHOR.MIDDLE, space=1.02)

text(s, 0.7, 6.85, 12.0, 0.5, [[R("Read:  ", 14, WEB3, True),
    R("being fundable means being a top name in your lane — the megaround dynamic is not AI-specific.", 14, MUTE, False, True)]])

# ============================================================ SLIDE 5 — the questions
s = slide()
kicker(s, "04  ·  WHAT TO WEIGH", AI)
text(s, 0.7, 0.78, 12.0, 0.7, [[R("Three questions the data puts in front of you", 30, INK, True, False, HEAD)]])

qs = [
    ("Attention vs. focus", AI,
     "AI holds ~50% of capital and the narrative.",
     "Do you need an AI story to be in the conversation — or does chasing one pull effort from where you can actually win?"),
    ("Depth in your lane", WEB3,
     "Crypto & fintech capital is back, but concentrated.",
     "Is your next unit of effort buying category leadership — or spreading you thin across features?"),
    ("Clearing the bar", FIN,
     "Megarounds mean a few names take most of the money.",
     "Where does effort move you closer to being the name that takes the round: AI features, or depth in your core?"),
]
cy, ch, cgap = 1.85, 1.5, 0.2
for i, (title, col, signal, q) in enumerate(qs):
    y = cy + i * (ch + cgap)
    card(s, 0.7, y, 11.95, ch)
    text(s, 1.05, y, 0.8, ch, [[R(str(i + 1), 30, col, True, False, HEAD)]], anchor=MSO_ANCHOR.MIDDLE)
    text(s, 1.95, y + 0.18, 10.4, ch - 0.3, [
        [R(title + "  —  ", 16, col, True), R(signal, 16, MUTE, False, True)],
        [R(q, 14.5, INK)],
    ], space=1.03, para_after=4, anchor=MSO_ANCHOR.MIDDLE)

text(s, 0.7, 6.95, 11.95, 0.45,
     [[R("The data sets the table. The call is yours.", 17, INK, True, True, HEAD)]],
     align=PP_ALIGN.CENTER)

# ---- speaker notes ----
notes = [
 "Framing slide. This is a founder's read on where capital went 2022-2025, not a market report. Every number is sourced (Crunchbase + Galaxy Digital). The audience question running underneath: how much should I lean into AI?",
 "The gravity. AI pulled ~50% of ALL global VC in 2025, $211B total (4.6x since 2022), and ~$93B of that went to just three labs. Point: the attention is real and overwhelming, but the AI core is concentrated and expensive to compete in directly. Don't tell them what to do - surface that the gravity exists.",
 "Their own lane. Crypto collapsed 86% then recovered to its strongest quarter since early 2022 - and into trading infrastructure, exactly the GTE category. Fintech halved then rebuilt on fewer, larger rounds. Both alive, both disciplined and concentrated.",
 "The cross-cutting pattern. Average landmark round grew ~6x to $2.6B. In every sector a few names absorb most of the money: AI half to 3 labs, Web3 41% to one deal, fintech 15% to Stripe. Fundability = being a top name in your lane.",
 "Synthesis - pose, don't prescribe. Three questions: do you need an AI story or does it pull focus; is effort buying category leadership; where does effort get you closer to taking the round. Close: the data sets the table, the call is theirs.",
]
for sl, nt in zip(prs.slides, notes):
    sl.notes_slide.notes_text_frame.text = nt

prs.save("/tmp/vcflow/Capital_Map_Founder_Brief.pptx")
print("saved")
